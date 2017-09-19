import asyncio
import math
import os
import sys
from base64 import b64encode
from datetime import datetime
from io import BytesIO

import aiohttp
import mercantile
from PIL import Image, ImageMath
from PIL import ImageDraw
from clover.geometry.bbox import BBox
from clover.render.renderers.stretched import StretchedRenderer
from clover.utilities.color import Color
from django.conf import settings
from django.contrib.gis.geos import Point
from django.template.loader import render_to_string
from geopy.distance import vincenty
from ncdjango.geoimage import world_to_image, image_to_world
from pyproj import Proj, transform
from weasyprint import HTML
from pptx import Presentation
from pptx.util import Inches, Pt

from seedsource.models import SeedZone, TransferLimit, Region
from seedsource.utils import get_elevation_at_point
from seedsource.report_config import VARIABLE_CONFIG, CONSTRAINT_CONFIG

ALLOWED_HOSTS = getattr(settings, 'ALLOWED_HOSTS')
BASE_DIR = settings.BASE_DIR
PORT = getattr(settings, 'PORT', 80)

TILE_SIZE = (256, 256)
IMAGE_SIZE = (645, 430)

SPECIES_LABELS = {
    'generic': 'Generic',
    'psme': 'Douglas-fir',
    'pico': 'Lodgepole pine',
    'piba': 'Jack pine',
    'pipo': 'Ponderosa pine',
    'pima': 'Black spruce',
    'thpl': 'Western redcedar',
    'pimo': 'Western white pine'
}

YEAR_LABELS = {
    '1961_1990': '1961-1990',
    '1981_2010': '1981-2010',
    '2025': '2011-2040',
    '2055': '2041-2070',
    '2085': '2071-2100'
}

RESULTS_RENDERER = StretchedRenderer([
    (0, Color(240, 59, 32)),
    (50, Color(254, 178, 76)),
    (100, Color(255, 237, 160))
])

DEGREE_SIGN = u'\N{DEGREE SIGN}'

# Map idxs
MAP_IDX = 15
SCALE_IDX = 21
WEST_IDX = 16
NORTH_IDX = 17
EAST_IDX = 18
SOUTH_IDX = 19
ATTRIBUTION_IDX = 20    # Copied throughout slide master, same in each slide

# Stats idxs
DATE_IDX = 27
CONFIG_TEXT_IDX = 34

# Table idxs
VAR_TABLE = 28

CONSTRAINTS = {
    'elevation': {
        'name': 'Elevation',
        'units': ('meters', 'feet'),
        'text': (('Between ', False),  ('{min}', True), (' and ', False), ('{max}', True), (' {units}', False))
    },
    'photoperiod': {
        'name': 'Photoperiod',
        'text': (('Within ', False), ('{hours:.1f}', True), (' hours of ', False), ('{month}/{day}/{year}', True))
    },
    'latitude': {
        'name': 'Latitude',
        'text': (('Between ', False), ('{min:.2f}', True), (' and ', False), ('{max:.2f}', True), (' &deg;N', False))
    },
    'longitude': {
        'name': 'Longitude',
        'text': (('Between ', False), ('{min:.2f}', True), (' and ', False), ('{max:.2f}', True), (' &deg;E', False))
    },
    'distance': {
        'name': 'Distance',
        'units': ('meters', 'feet'),
        'text': (('Within ', False), ('{distance:.1f} {units}', True))
    }
}

ATTRIBUTION_TEXT = 'Generated by the Seedlot Selection Tool'


def degree_sign(string):
    return string.replace('&deg;', DEGREE_SIGN)


class Report(object):
    def __init__(self, configuration, zoom, tile_layers, opacity):
        self.configuration = configuration
        self.zoom = zoom
        self.tile_layers = tile_layers
        self.opacity = opacity

    def get_year(self, climate):
        return YEAR_LABELS[climate['time']]

    def get_model(self, climate):
        if climate['time'] in {'1961_1990', '1981_2010'}:
            return None
        else:
            return {
                'rcp45': 'RCP4.5',
                'rcp85': 'RCP8.5'
            }[climate['model']]

    def get_context_variables(self):
        variables = []
        is_imperial = self.configuration['unit'] == 'imperial'

        for variable in self.configuration['variables']:
            name, value, transfer = variable['name'], variable['value'], variable['transfer']
            config = VARIABLE_CONFIG[name]
            value /= config.multiplier
            transfer /= config.multiplier

            variables.append({
                'label': '{}: {}'.format(variable['name'], config.label),
                'value': config.format_value(value, is_imperial),
                'limit': config.format_transfer(transfer, is_imperial),
                'units': config.imperial_label if is_imperial else config.metric_label,
                'modified': variable['transfer'] != variable['defaultTransfer']
            })

        return variables

    def get_context_constraints(self):
        constraints = []
        is_imperial = self.configuration['unit'] == 'imperial'

        for constraint in self.configuration['constraints']:
            name, values = constraint['type'], constraint['values']
            config = CONSTRAINT_CONFIG[name]

            constraints.append({
                'label': config.label,
                'value': config.format_value(self.configuration, is_imperial),
                'range': config.format_range(values, is_imperial)
            })

        return constraints

    def get_context(self, img_as_bytes=False):
        point = self.configuration['point']
        elevation = get_elevation_at_point(Point(point['x'], point['y'])) / 0.3048
        method = self.configuration['method']
        objective = self.configuration['objective']
        climates = self.configuration['climate']

        if self.configuration['method'] == 'seedzone' and elevation is not None:
            zone_uid = self.configuration['zones']['selected']

            try:
                zone = SeedZone.objects.get(zone_uid=zone_uid)
                zone_id = zone.pk

                try:
                    limit = zone.transferlimit_set.filter(low__lt=elevation, high__gte=elevation)[:1].get()
                    band = [0 if limit.low == -1 else limit.low, limit.high]
                except TransferLimit.DoesNotExist:
                    band = None
            except SeedZone.DoesNotExist:
                zone_id = None
                zone = None
                band = None
        else:
            zone_id = None
            zone = None
            band = None

        mercator = Proj(init='epsg:3857')
        wgs84 = Proj(init='epsg:4326')

        map_image, map_bbox = MapImage(
            IMAGE_SIZE, (point['x'], point['y']), self.zoom, self.tile_layers, self.configuration.get('region'),
            zone_id, self.opacity
        ).get_image()
        to_world = image_to_world(map_bbox, map_image.size)
        map_bbox = map_bbox.project(Proj(init='epsg:4326'), edge_points=0)

        image_data = BytesIO()
        map_image.save(image_data, 'png')
        if img_as_bytes:
            image_data.seek(0)  # seek to file start so we can read into ppt

        with open(os.path.join(BASE_DIR, 'seedsource', 'static', 'sst', 'images', 'north.png'), 'rb') as f:
            north_image_data = b64encode(f.read())

        with open(os.path.join(BASE_DIR, 'seedsource', 'static', 'sst', 'images', 'scale.png'), 'rb') as f:
            scale_image_data = b64encode(f.read())

        scale_bar_x = 38
        scale_bar_y = map_image.size[1] - 15
        scale_bar_start = transform(mercator, wgs84, *to_world(scale_bar_x, scale_bar_y))
        scale_bar_end = transform(mercator, wgs84, *to_world(scale_bar_x + 96, scale_bar_y))
        scale = '{} mi'.format(round(vincenty(scale_bar_start, scale_bar_end).miles, 1))

        legend = RESULTS_RENDERER.get_legend()[0]

        def format_x_coord(x):
            return '{}&deg; W'.format(round(abs(x), 2)) if x < 0 else '{}&deg; E'.format(round(x, 2))

        def format_y_coord(y):
            return '{}&deg; S'.format(round(abs(y), 2)) if y < 0 else '{}&deg; N'.format(round(y, 2))

        return {
            'today': datetime.today(),
            'image_data': b64encode(image_data.getvalue()) if not img_as_bytes else image_data,
            'north': format_y_coord(map_bbox.ymax),
            'east': format_x_coord(map_bbox.xmax),
            'south': format_y_coord(map_bbox.ymin),
            'west': format_x_coord(map_bbox.xmin),
            'north_image_data': north_image_data,
            'scale_image_data': scale_image_data,
            'scale': scale,
            'legend_image_data': legend.image_base64,
            'objective': 'Find seedlots' if objective == 'seedlots' else 'Find planting sites',
            'location_label': 'Planting site location' if objective == 'seedlots' else 'Seedlot location',
            'point': {'x': round(point['x'], 4), 'y': round(point['y'], 4)},
            'elevation': round(elevation),
            'seedlot_year': self.get_year(climates['seedlot']),
            'site_year': self.get_year(climates['site']),
            'site_model': self.get_model(climates['site']),
            'method': method,
            'center': self.configuration['center'],
            'species': SPECIES_LABELS[self.configuration['species']] if method == 'seedzone' else None,
            'zone': getattr(zone, 'name', None),
            'band': band,
            'variables': self.get_context_variables(),
            'constraints': self.get_context_constraints()
        }

    def get_pdf_data(self) -> BytesIO:
        pdf_data = BytesIO()

        HTML(
            BytesIO(render_to_string('pdf/report.html', self.get_context()).encode())
        ).write_pdf(pdf_data)

        return pdf_data

    def get_pptx_data(self) -> BytesIO:
        # Add a text frame where needed
        def add_text_frame(tframe, lines):
            p = tframe.paragraphs[0]
            for line in lines:
                bold_text, norm_text, add_newline = line
                run = p.add_run()
                run.text = bold_text
                run.font.bold = True
                run = p.add_run()
                run.text = norm_text + '\n'
                if add_newline:
                    run = p.add_run()
                    run.text = '\n'

        ppt_data = BytesIO()

        ctx = self.get_context(img_as_bytes=True)
        pptx_template_path = os.path.join(settings.BASE_DIR, 'seedsource', 'static', 'sst', 'ppt',
                                          'report_template7.pptx')
        prs = Presentation(pptx_template_path)

        t = ctx['today']
        date = t.strftime('%m/%d/%Y')
        attribution = '{} on {}'.format(ATTRIBUTION_TEXT, date)

        # Set map slide placeholders
        mapslide = prs.slides[0]
        placeholder = mapslide.placeholders[MAP_IDX]
        placeholder.insert_picture(ctx['image_data'])
        placeholder = mapslide.placeholders[SCALE_IDX]
        placeholder.text = ctx['scale']
        placeholder = mapslide.placeholders[WEST_IDX]
        placeholder.text = degree_sign(ctx['west'])
        placeholder = mapslide.placeholders[NORTH_IDX]
        placeholder.text = degree_sign(ctx['north'])
        placeholder = mapslide.placeholders[EAST_IDX]
        placeholder.text = degree_sign(ctx['east'])
        placeholder = mapslide.placeholders[SOUTH_IDX]
        placeholder.text = degree_sign(ctx['south'])
        placeholder = mapslide.placeholders[ATTRIBUTION_IDX]
        placeholder.text = attribution

        # Gather all text for mapslide.notes_slide with all the stat info used below and in statslide
        objective = ctx['objective']
        location_label = ctx['location_label'] + ':'
        point = ctx['point']
        point_label = 'Lat: {y}{degree_sign}, Lon: {x}{degree_sign}'.format(x=point['x'], y=point['y'],
                                                                            degree_sign=DEGREE_SIGN)
        elev_label = '{} ft'.format(ctx['elevation'])
        seedlot_year = ctx['seedlot_year']
        site_year = ctx['site_year']
        tmethod = ctx['method']
        center = ctx['center']

        # Determine method text
        if tmethod == 'seedzone':
            if center == 'zone':
                method_text = 'Transfer limits and climatic center based on seed zone'
            else:
                method_text = 'Transfer limits based on seed zone, climatic center based on the selected location'
        else:
            method_text = 'Custom transfer limits, climatic center based on the selected location'

        lines = [
            ('Objective: ', objective, True),               # (bold_text, norm_text, add_newline)
            (location_label + ' ', point_label, False),
            ('Elevation: ', elev_label, True),
            ('Climate Scenarios', '', False),
            ('', 'Seedlot climate: ' + seedlot_year, False),
            ('', 'Planting site climate: ' + site_year, True),
            ('Transfer limit method: ', method_text, True)
        ]

        if tmethod == 'seedzone':
            species_text = 'Species: '
            species_label = ctx['species']
            sz_text = 'Seed zone:'
            band_str = ", {}' - {}'".format(ctx['band'][0], ctx['band'][1]) if ctx['band'] else ''
            sz_label = '{}{}'.format(ctx['zone'], band_str)
            lines += [
                (species_text, species_label, False),
                (sz_text, sz_label, False)
            ]

        name_width = max([len('Variable')] + [len(x['label']) for x in ctx['variables']]) + 3
        center_width = max(
            [len('Center')] + [len(' '.join([str(x['value']), degree_sign(x['units'])])) for x in ctx['variables']]
        ) + 3
        transfer_width = max(
            [len('Transfer limit (+/-)')] +
            [
                len('{} {}{}'.format(x['limit'], degree_sign(x['units']), ' (modified)' if x['modified'] else ''))
                for x in ctx['variables']
            ]
        )

        variable_notes = [
            ('Variables', '', False),
            ('', ''.join([
                'Variable'.ljust(name_width),
                'Center'.ljust(center_width),
                'Transfer limit (+/-)'.ljust(transfer_width)
            ]), False),
            ('', '-' * (name_width + center_width + transfer_width), False)
        ]

        for variable in ctx['variables']:
            units = degree_sign(variable['units'])
            variable_notes.append(('', ''.join([
                variable['label'].ljust(name_width),
                '{} {}'.format(variable['value'], units).ljust(center_width),
                '{} {}{}'.format(variable['limit'], units, ' (modified)' if variable['modified'] else '')
            ]), False))

        variable_notes.append(('', '', False))

        notes_slide = mapslide.notes_slide
        notes_slide.notes_text_frame.paragraphs[0].font.name = 'Andale Mono'
        add_text_frame(notes_slide.notes_text_frame, lines + variable_notes)

        # Set run config placeholders
        statslide = prs.slides[1]
        placeholder = statslide.placeholders[ATTRIBUTION_IDX]
        placeholder.text = attribution
        placeholder = statslide.placeholders[DATE_IDX]
        placeholder.text = date
        placeholder = statslide.placeholders[CONFIG_TEXT_IDX]
        add_text_frame(placeholder.text_frame, lines)

        # Populate variables slide
        variable_slide = prs.slides[2]
        placeholder = variable_slide.placeholders[ATTRIBUTION_IDX]
        placeholder.text = attribution

        # Insert table
        placeholder = variable_slide.placeholders[VAR_TABLE]
        numrows = len(ctx['variables'])+1  # +1 row for column headings
        numcols = 3
        table = placeholder.insert_table(rows=numrows, cols=numcols).table
        table.horz_banding = False
        table.vert_banding = False

        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(3.0)

        # Fill column headings
        table.cell(0, 0).text = 'Variable'
        table.cell(0, 1).text = 'Center'
        table.cell(0, 2).text = 'Transfer limit (+/-)'

        # Now fill cells
        for row, variable in enumerate(ctx['variables'], start=1):
            units = degree_sign(variable['units'])
            center_label = '{} {}'.format(variable['value'], units)
            limit_label = '{} {}{}'.format(variable['limit'], units, ' (modified)' if variable['modified'] else '')
            table.cell(row, 0).text = variable['label']
            table.cell(row, 1).text = center_label
            table.cell(row, 2).text = limit_label

        # Now format the text. Easier to do this after being set
        for i in range(numrows):
            font_size = Pt(18) if i == 0 else Pt(16)
            for j in range(numcols):
                p = table.cell(i, j).text_frame.paragraphs[0]
                p.font.size = font_size

        if ctx['constraints']:
            slide = prs.slides.add_slide(prs.slide_layouts[3])

            title_tx = slide.shapes.add_textbox(378320, 393424, 2090559, 369332)
            tf = title_tx.text_frame
            tf.text = 'Constraints'
            tf.paragraphs[0].font.bold = True

            placeholder = variable_slide.placeholders[ATTRIBUTION_IDX]
            attribution_tx = slide.shapes.add_textbox(
                placeholder.left, placeholder.top, placeholder.width, placeholder.height
            )
            tf = attribution_tx.text_frame
            tf.text = attribution
            tf.paragraphs[0].font.size = Pt(12)

            constraints_tx = slide.shapes.add_textbox(378320, Inches(1), Inches(8), Inches(2))
            tf = constraints_tx.text_frame

            notes_label = notes_slide.notes_text_frame.paragraphs[0].add_run()
            notes_label.font.bold = True
            notes_label.text = '\nConstraints\n'

            for paragraph in (tf.paragraphs[0], notes_slide.notes_text_frame.paragraphs[0]):
                for constraint in self.configuration['constraints']:
                    info = CONSTRAINTS[constraint['name']]
                    run = paragraph.add_run()
                    run.font.bold = True
                    run.text = '{}: '.format(info['name'])

                    for text in info['text']:
                        run = paragraph.add_run()
                        run.font.bold = text[1]
                        run.text = degree_sign(text[0].format(**constraint['args']))

                    paragraph.add_run().text = '\n'


        prs.save(ppt_data)
        ppt_data.seek(0)
        return ppt_data


class MapImage(object):
    def __init__(self, size, point, zoom, tile_layers, region, zone_id, opacity):
        self._configure_event_loop()

        self.num_tiles = [math.ceil(size[x] / TILE_SIZE[x]) + 1 for x in (0, 1)]
        center_tile = mercantile.tile(point[0], point[1], zoom)

        mercator = Proj(init='epsg:3857')
        wgs84 = Proj(init='epsg:4326')

        center_tile_bbox = BBox(mercantile.bounds(*center_tile), projection=wgs84).project(mercator, edge_points=0)
        center_to_image = world_to_image(center_tile_bbox, TILE_SIZE)
        center_to_world = image_to_world(center_tile_bbox, TILE_SIZE)
        center_point_px = center_to_image(*mercantile.xy(*point))

        self.ul_tile = mercantile.tile(
            *transform(mercator, wgs84, *center_to_world(
                center_point_px[0] - math.ceil(IMAGE_SIZE[0] / 2),
                center_point_px[1] - math.ceil(IMAGE_SIZE[1] / 2)
            ), zoom))

        lr_tile = mercantile.Tile(
            x=min(2 ** zoom, self.ul_tile.x + self.num_tiles[0]),
            y=min(2 ** zoom, self.ul_tile.y + self.num_tiles[1]),
            z=zoom
        )

        ul = mercantile.xy(*mercantile.ul(*self.ul_tile))
        lr = mercantile.xy(*mercantile.ul(*lr_tile))

        self.image_bbox = BBox((ul[0], lr[1], lr[0], ul[1]))
        self.image_size = (TILE_SIZE[0] * self.num_tiles[0], TILE_SIZE[1] * self.num_tiles[1])

        self.to_image = world_to_image(self.image_bbox, self.image_size)
        self.to_world = image_to_world(self.image_bbox, self.image_size)

        self.point_px = [round(x) for x in self.to_image(*mercantile.xy(*point))]

        self.target_size = size
        self.point = point
        self.zoom = zoom
        self.tile_layers = tile_layers
        self.region = region
        self.zone_id = zone_id
        self.opacity = opacity

    def _configure_event_loop(self):
        if sys.platform == 'win32':
            asyncio.set_event_loop(asyncio.ProactorEventLoop())
        else:
            asyncio.set_event_loop(asyncio.SelectorEventLoop())

    def get_layer_images(self):
        async def fetch_tile(client, layer_url, tile, im):
            headers = {}

            layer_url = layer_url.format(x=tile.x, y=tile.y, z=tile.z, s='server')
            if layer_url.startswith('//'):
                layer_url = 'https:{}'.format(layer_url)
            elif layer_url.startswith('/'):
                layer_url = 'http://127.0.0.1:{}{}'.format(PORT, layer_url)
                if ALLOWED_HOSTS:
                    headers['Host'] = ALLOWED_HOSTS[0]

            async with client.get(layer_url, headers=headers) as r:
                tile_im = Image.open(BytesIO(await r.read()))
                im.paste(tile_im, ((tile.x - self.ul_tile.x) * 256, (tile.y - self.ul_tile.y) * 256))

        layer_images = [Image.new('RGBA', self.image_size) for _ in self.tile_layers]

        with aiohttp.ClientSession() as client:
            requests = []

            for i in range(self.num_tiles[0] * self.num_tiles[1]):
                tile = mercantile.Tile(
                    x=self.ul_tile.x + i % self.num_tiles[0],
                    y=self.ul_tile.y + i // self.num_tiles[0],
                    z=self.zoom
                )

                for j, layer_url in enumerate(self.tile_layers):
                    requests.append(fetch_tile(client, layer_url, tile, layer_images[j]))

            asyncio.get_event_loop().run_until_complete(asyncio.gather(*requests))

        return layer_images

    def draw_geometry(self, im, geometry, color, width):
        canvas = ImageDraw.Draw(im)

        canvas.line(
            [tuple(round(x) for x in self.to_image(*mercantile.xy(*p))) for p in geometry], fill=color, width=width
        )

    def draw_zone_geometry(self, im):
        if self.zone_id is not None:
            polygon = SeedZone.objects.get(pk=self.zone_id).polygon

            if polygon.geom_type == 'MultiPolygon':
                geometries = polygon.coords
            else:
                geometries = [polygon.coords]

            for geometry in geometries:
                self.draw_geometry(im, geometry[0], (0, 255, 0), 3)

    def draw_region_geometry(self, im):
        try:
            region = Region.objects.filter(name=self.region).get()
        except Region.DoesNotExist:
            return

        for geometry in region.polygons.coords:
            self.draw_geometry(im, geometry[0], (0, 0, 102), 1)

    def get_marker_image(self):
        leaflet_images_dir = os.path.join(BASE_DIR, 'seedsource', 'static', 'leaflet', 'images')
        marker = Image.open(os.path.join(leaflet_images_dir, 'marker-icon.png'))
        shadow = Image.open(os.path.join(leaflet_images_dir, 'marker-shadow.png'))

        # Raise the shadow opacity
        shadow.putalpha(ImageMath.eval('a * 2', a=shadow.convert('RGBA').split()[3]).convert('L'))

        im = Image.new('RGBA', self.image_size)
        im.paste(shadow, (self.point_px[0] - 12, self.point_px[1] - shadow.size[1]))

        marker_im = Image.new('RGBA', im.size)
        marker_im.paste(marker, (self.point_px[0] - marker.size[0] // 2, self.point_px[1] - marker.size[1]))
        im.paste(marker_im, (0, 0), marker_im)

        return im

    def crop_image(self, im):
        im_ul = (self.point_px[0] - self.target_size[0] // 2, self.point_px[1] - self.target_size[1] // 2)
        box = (*im_ul, im_ul[0] + self.target_size[0], im_ul[1] + self.target_size[1])

        return im.crop(box), BBox(
            (self.to_world(box[0], box[3])) + self.to_world(box[2], box[1]), projection=Proj(init='epsg:3857')
        )

    def get_image(self) -> (Image, BBox):
        im = Image.new('RGBA', self.image_size)

        for i, layer_im in enumerate(self.get_layer_images()):
            im.paste(Image.blend(im, layer_im, 1 if i == 0 else self.opacity), (0, 0), layer_im)

        self.draw_zone_geometry(im)
        self.draw_region_geometry(im)

        marker_im = self.get_marker_image()
        im.paste(marker_im, (0, 0), marker_im)

        return self.crop_image(im)
